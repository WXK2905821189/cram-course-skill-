import importlib.util
import json
import tempfile
import unittest
from datetime import date
from pathlib import Path


ROOT = Path(__file__).resolve().parents[1]


def load_module(name, relative_path):
    module_path = ROOT / relative_path
    spec = importlib.util.spec_from_file_location(name, module_path)
    module = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(module)
    return module


extract_content = load_module("extract_content", "cram-course/scripts/extract_content.py")
spaced_repetition = load_module("spaced_repetition", "cram-course/scripts/spaced_repetition.py")


class ExtractContentTests(unittest.TestCase):
    def test_identify_knowledge_points_extracts_titles_and_definitions(self):
        result = extract_content.identify_knowledge_points([
            {
                "id": 1,
                "title": "栈的基本概念",
                "content": [
                    "栈是只允许在一端进行插入和删除的线性表。",
                    "S = top - base"
                ],
                "source": "sample.pptx"
            }
        ])

        self.assertEqual(result["knowledge_points"][0]["title"], "栈的基本概念")
        self.assertTrue(any("栈是" in item["text"] for item in result["definitions"]))
        self.assertTrue(any("S =" in item["text"] for item in result["formulas"]))

    def test_extract_from_pptx_handles_blank_slide_and_titles(self):
        from pptx import Presentation

        with tempfile.TemporaryDirectory() as tmpdir:
            pptx_path = Path(tmpdir) / "sample.pptx"
            prs = Presentation()
            title_slide = prs.slides.add_slide(prs.slide_layouts[0])
            title_slide.shapes.title.text = "栈与队列"
            title_slide.placeholders[1].text = "栈是后进先出的线性表。"
            prs.slides.add_slide(prs.slide_layouts[6])
            prs.save(pptx_path)

            slides = extract_content.extract_from_pptx(str(pptx_path))

        self.assertEqual(len(slides), 2)
        self.assertEqual(slides[0]["title"], "栈与队列")
        self.assertEqual(slides[1]["title"], "幻灯片 2")

    def test_extract_from_docx_preserves_headings_and_tables(self):
        from docx import Document

        with tempfile.TemporaryDirectory() as tmpdir:
            docx_path = Path(tmpdir) / "sample.docx"
            doc = Document()
            doc.add_heading("第一章 栈", level=1)
            doc.add_paragraph("栈是只允许在一端插入和删除的线性表。")
            table = doc.add_table(rows=1, cols=2)
            table.cell(0, 0).text = "结构"
            table.cell(0, 1).text = "特点"
            doc.save(docx_path)

            result = extract_content.extract_from_docx(str(docx_path))

        self.assertEqual(result["sections"][0]["title"], "第一章 栈")
        self.assertIn("栈是", result["sections"][0]["content"][0])
        self.assertEqual(result["tables"][0]["data"][0], ["结构", "特点"])


class SpacedRepetitionTests(unittest.TestCase):
    def test_parse_date_accepts_common_formats(self):
        self.assertEqual(spaced_repetition.parse_date("2026-07-15"), date(2026, 7, 15))
        self.assertEqual(spaced_repetition.parse_date("2026/07/15"), date(2026, 7, 15))

    def test_load_topics_deduplicates_titles(self):
        data = {
            "knowledge_points": [
                {"title": "栈", "source": "a.pptx"},
                {"title": "栈", "source": "a.pptx"}
            ],
            "raw_content": [
                {"title": "队列", "source": "a.pptx"},
                {"title": "幻灯片 3", "source": "a.pptx"}
            ]
        }

        with tempfile.NamedTemporaryFile("w", encoding="utf-8", suffix=".json", delete=False) as tmp:
            json.dump(data, tmp, ensure_ascii=False)
            tmp_path = tmp.name

        try:
            topics = spaced_repetition.load_topics(tmp_path)
        finally:
            Path(tmp_path).unlink(missing_ok=True)

        self.assertEqual([topic["title"] for topic in topics], ["栈", "队列"])


class RepositoryPackagingTests(unittest.TestCase):
    def test_license_exists_for_github_release(self):
        self.assertTrue((ROOT / "LICENSE").exists())

    def test_manifest_declares_supported_formats_and_outputs(self):
        manifest_path = ROOT / "manifest.json"
        self.assertTrue(manifest_path.exists())
        manifest = json.loads(manifest_path.read_text(encoding="utf-8"))
        self.assertEqual(manifest["name"], "cram-course")
        self.assertIn("pptx", manifest["supported_inputs"])
        self.assertIn("02_覆盖矩阵.md", manifest["required_outputs"])
        self.assertIn("30_名词解释专项.md", manifest["required_outputs"])
        self.assertNotIn("coverage_matrix.md", manifest["required_outputs"])
        self.assertNotIn("mock_exam.md", manifest["required_outputs"])

    def test_coverage_matrix_is_required_and_templated(self):
        skill_text = (ROOT / "cram-course/SKILL.md").read_text(encoding="utf-8-sig")
        self.assertIn("02_覆盖矩阵.md", skill_text)
        self.assertTrue((ROOT / "cram-course/assets/coverage_matrix_template.md").exists())

    def test_subject_modes_and_material_diagnostics_are_routed(self):
        skill_text = (ROOT / "cram-course/SKILL.md").read_text(encoding="utf-8-sig")
        self.assertIn("references/subject_modes.md", skill_text)
        self.assertIn("references/material_diagnostics.md", skill_text)
        self.assertTrue((ROOT / "cram-course/references/subject_modes.md").exists())
        self.assertTrue((ROOT / "cram-course/references/material_diagnostics.md").exists())

    def test_full_demo_contains_core_artifacts(self):
        demo = ROOT / "examples/full-demo"
        self.assertTrue((demo / "README.md").exists())
        self.assertTrue((demo / "02_覆盖矩阵.md").exists())
        self.assertTrue((demo / "01_学习计划.md").exists())
        self.assertTrue((demo / "10_第01章_数据结构基础_02章节题库.md").exists())

    def test_user_facing_markdown_names_are_chinese_and_obsidian_ordered(self):
        skill_text = (ROOT / "cram-course/SKILL.md").read_text(encoding="utf-8-sig")
        expected_names = [
            "00_课程导学.md",
            "01_学习计划.md",
            "02_覆盖矩阵.md",
            "10_第01章_<章节名>_01学习笔记.md",
            "10_第01章_<章节名>_02章节题库.md",
            "30_名词解释专项.md",
            "40_间隔复习计划.md",
            "50_模拟考试.md",
        ]
        for name in expected_names:
            with self.subTest(name=name):
                self.assertIn(name, skill_text)

    def test_templates_warn_ai_outputs_need_review(self):
        template_paths = [
            ROOT / "cram-course/assets/study_plan_template.md",
            ROOT / "cram-course/assets/note_template.md",
            ROOT / "cram-course/assets/topic_drill_template.md",
            ROOT / "cram-course/assets/error_log_template.md",
            ROOT / "cram-course/assets/coverage_matrix_template.md",
            ROOT / "cram-course/references/question_templates.md",
        ]
        for path in template_paths:
            with self.subTest(path=path.name):
                text = path.read_text(encoding="utf-8-sig")
                self.assertIn("人工复核", text)

    def test_skill_entrypoint_stays_concise(self):
        skill_lines = (ROOT / "cram-course/SKILL.md").read_text(encoding="utf-8-sig").splitlines()
        self.assertLessEqual(len(skill_lines), 260)


if __name__ == "__main__":
    unittest.main()
