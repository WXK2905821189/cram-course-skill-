#!/usr/bin/env python3
"""
extract_content.py — 从 PPT/PDF/DOCX 课件中提取结构化内容

用法:
    python extract_content.py <file1> <file2> ... --output study_materials.json

输出: JSON 文件，包含章节结构、知识点、定义、公式等。
支持 .pptx/.pdf/.docx；旧版 .ppt 请先转换为 .pptx。
"""

import argparse
import json
import os
import re
import sys
from pathlib import Path


def configure_utf8_output():
    """Keep CLI output readable in Git Bash, CI, and modern Windows terminals."""
    for stream in (sys.stdout, sys.stderr):
        if hasattr(stream, "reconfigure"):
            stream.reconfigure(encoding="utf-8")


def extract_from_pptx(filepath):
    """从 PPTX 文件中提取文本，按幻灯片组织。"""
    try:
        from pptx import Presentation
    except ImportError:
        print("错误: 需要安装 python-pptx。请运行: pip install python-pptx")
        sys.exit(1)

    prs = Presentation(filepath)
    slides_data = []
    current_section = ""

    for i, slide in enumerate(prs.slides, 1):
        texts = []
        title = f"幻灯片 {i}"

        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        texts.append(text)

            if shape.has_table:
                table = shape.table
                table_data = []
                for row in table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    table_data.append(row_data)
                texts.append({"type": "table", "data": table_data})

        # 尝试检测章节标题（第一个非空文本）
        string_texts = [text for text in texts if isinstance(text, str)]
        slide_title = string_texts[0] if string_texts else title
        if slide_title and len(slide_title) < 80 and not slide_title.startswith("•"):
            title = slide_title

        slides_data.append({
            "id": i,
            "title": title,
            "content": texts,
            "source": os.path.basename(filepath)
        })

    return slides_data


def extract_from_pdf(filepath):
    """从 PDF 文件中提取文本。"""
    try:
        from PyPDF2 import PdfReader
    except ImportError:
        print("错误: 需要安装 PyPDF2。请运行: pip install PyPDF2")
        sys.exit(1)

    reader = PdfReader(filepath)
    pages_data = []
    current_page = 1

    for page in reader.pages:
        text = page.extract_text()
        if text and text.strip():
            lines = [line.strip() for line in text.split("\n") if line.strip()]
            # 尝试识别标题（短行、可能全部大写或加粗特征）
            title = lines[0] if lines else f"第 {current_page} 页"
            if len(title) > 100:
                title = f"第 {current_page} 页"

            pages_data.append({
                "id": current_page,
                "title": title,
                "content": lines,
                "source": os.path.basename(filepath)
            })
        current_page += 1

    return pages_data


def extract_from_docx(filepath):
    """从 DOCX 文件中提取文本，按段落和标题组织。"""
    try:
        from docx import Document
    except ImportError:
        print("错误: 需要安装 python-docx。请运行: pip install python-docx")
        sys.exit(1)

    doc = Document(filepath)
    sections = []
    current_section = None

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue

        # 检测标题样式
        if para.style.name.startswith("Heading"):
            if current_section:
                sections.append(current_section)
            current_section = {
                "title": text,
                "level": int(para.style.name.split()[-1]) if para.style.name.split()[-1].isdigit() else 1,
                "content": [],
                "source": os.path.basename(filepath)
            }
        elif current_section is not None:
            current_section["content"].append(text)
        else:
            # 没有标题的段落归入默认节
            if not sections:
                current_section = {
                    "title": "正文",
                    "level": 1,
                    "content": [],
                    "source": os.path.basename(filepath)
                }
            current_section["content"].append(text)

    if current_section:
        sections.append(current_section)

    # 提取表格
    tables_data = []
    for i, table in enumerate(doc.tables):
        table_content = []
        for row in table.rows:
            row_data = [cell.text.strip() for cell in row.cells]
            table_content.append(row_data)
        tables_data.append({
            "id": i + 1,
            "type": "table",
            "data": table_content
        })

    return {"sections": sections, "tables": tables_data, "source": os.path.basename(filepath)}


def identify_knowledge_points(all_content):
    """从提取的内容中识别知识点、定义和公式。"""
    knowledge_points = []
    definitions = []
    formulas = []

    # 定义检测模式
    def_patterns = [
        r"(.{1,30}(?:是|指|即|定义为|是指|指的是).{5,200}?[。；;.])",
        r"(.{1,30}?:.{5,200}?[。；;.])",
    ]

    # 公式检测模式 (包含数学符号或 LaTeX 标记)
    formula_patterns = [
        r"([A-Za-z]\s*=\s*[^。；;]{3,100})",  # X = ...
        r"(\$[^$]+\$)",  # LaTeX inline
    ]

    for item in all_content:
        content_text = ""
        if isinstance(item, dict):
            if "content" in item:
                if isinstance(item["content"], list):
                    content_text = " ".join(
                        c if isinstance(c, str) else ""
                        for c in item["content"]
                    )
        elif isinstance(item, str):
            content_text = item

        if not content_text:
            continue

        # 提取定义
        for pattern in def_patterns:
            matches = re.findall(pattern, content_text)
            for m in matches:
                definitions.append({"text": m.strip(), "source": item.get("source", "unknown")})

        # 提取公式
        for pattern in formula_patterns:
            matches = re.findall(pattern, content_text)
            for m in matches:
                formulas.append({"text": m.strip(), "source": item.get("source", "unknown")})

        # 知识点：提取标题和要点
        if isinstance(item, dict) and item.get("title") and item["title"] != f"幻灯片 {item.get('id', '')}" and item["title"] != f"第 {item.get('id', '')} 页" and item["title"] != "正文":
            knowledge_points.append({
                "title": item["title"],
                "source": item.get("source", "unknown"),
                "type": "section"
            })

    return {
        "knowledge_points": knowledge_points,
        "definitions": definitions,
        "formulas": formulas
    }


def main():
    configure_utf8_output()

    parser = argparse.ArgumentParser(description="从 PPT/PDF/DOCX 课件提取结构化内容")
    parser.add_argument("files", nargs="+", help="课件文件路径 (支持 .pptx/.pdf/.docx；.ppt 请先转换为 .pptx)")
    parser.add_argument("--output", "-o", default="study_materials.json", help="输出 JSON 文件路径")
    args = parser.parse_args()

    all_content = []
    file_summary = []

    for filepath in args.files:
        path = Path(filepath)
        if not path.exists():
            print(f"警告: 文件不存在，跳过: {filepath}")
            continue

        ext = path.suffix.lower()
        print(f"正在处理: {filepath} ({ext})")

        try:
            if ext == ".pptx":
                content = extract_from_pptx(str(path))
                all_content.extend(content)
                file_summary.append({"file": str(path), "type": "pptx", "slides": len(content)})

            elif ext == ".ppt":
                print(f"警告: 暂不直接支持旧版 .ppt，请先转换为 .pptx: {filepath}")
                continue

            elif ext == ".pdf":
                content = extract_from_pdf(str(path))
                all_content.extend(content)
                file_summary.append({"file": str(path), "type": "pdf", "pages": len(content)})

            elif ext == ".docx":
                result = extract_from_docx(str(path))
                sections = result.get("sections", [])
                tables = result.get("tables", [])
                all_content.extend(sections)
                all_content.append({"tables": tables, "source": str(path)})
                file_summary.append({
                    "file": str(path),
                    "type": "docx",
                    "sections": len(sections),
                    "tables": len(tables)
                })

            else:
                print(f"警告: 不支持的文件类型: {ext}")
                continue

        except Exception as e:
            print(f"错误: 处理 {filepath} 时出错: {e}")
            continue

    # 识别知识点
    extracted = identify_knowledge_points(all_content)

    # 组装输出
    output = {
        "files": file_summary,
        "total_items": len(all_content),
        "raw_content": all_content,
        "knowledge_points": extracted["knowledge_points"],
        "definitions": extracted["definitions"],
        "formulas": extracted["formulas"]
    }

    # 写入 JSON
    with open(args.output, "w", encoding="utf-8") as f:
        json.dump(output, f, ensure_ascii=False, indent=2)

    print(f"\n[OK] 提取完成!")
    print(f"   - 处理文件: {len(file_summary)} 个")
    print(f"   - 提取条目: {len(all_content)} 条")
    print(f"   - 识别知识点: {len(extracted['knowledge_points'])} 个")
    print(f"   - 识别定义: {len(extracted['definitions'])} 条")
    print(f"   - 识别公式: {len(extracted['formulas'])} 条")
    print(f"   - 输出文件: {args.output}")


if __name__ == "__main__":
    main()
